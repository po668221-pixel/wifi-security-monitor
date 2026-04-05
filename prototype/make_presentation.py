from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colour palette ──────────────────────────────────────────────
BG       = RGBColor(0x0a, 0x0e, 0x1a)   # dark navy
CARD     = RGBColor(0x1a, 0x22, 0x35)
ACCENT   = RGBColor(0x00, 0xe5, 0xff)   # cyan
CRITICAL = RGBColor(0xef, 0x44, 0x44)
HIGH     = RGBColor(0xf9, 0x73, 0x16)
MEDIUM   = RGBColor(0xea, 0xb3, 0x08)
WHITE    = RGBColor(0xe2, 0xe8, 0xf0)
MUTED    = RGBColor(0x94, 0xa3, 0xb8)
DIM      = RGBColor(0x4b, 0x56, 0x63)

BLANK = prs.slide_layouts[6]   # truly blank layout


# ── Helpers ────────────────────────────────────────────────────

def bg(slide):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def box(slide, l, t, w, h, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return tb

def label_card(slide, title, body_lines, l, t, w, h, title_color=ACCENT):
    box(slide, l, t, w, h, fill_color=CARD)
    txt(slide, title, l+0.18, t+0.15, w-0.3, 0.3, size=10, bold=True, color=title_color)
    body = "\n".join(body_lines)
    txt(slide, body, l+0.18, t+0.5, w-0.3, h-0.65, size=13, color=WHITE)

def number_badge(slide, number, l, t, color):
    box(slide, l, t, 0.45, 0.45, fill_color=color)
    txt(slide, str(number), l, t, 0.45, 0.45, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def step_row(slide, num, heading, body, l, t, num_color=ACCENT):
    number_badge(slide, num, l, t, num_color)
    txt(slide, heading, l+0.6, t, 9, 0.28, size=14, bold=True, color=WHITE)
    txt(slide, body,    l+0.6, t+0.28, 9, 0.35, size=12, color=MUTED)

def emoji_bullet(slide, emoji, text, l, t, w=11):
    txt(slide, emoji, l, t, 0.5, 0.35, size=18)
    txt(slide, text,  l+0.55, t, w-0.55, 0.35, size=13, color=WHITE)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

# glow strip
box(s, 0, 2.8, 13.33, 0.04, fill_color=ACCENT)

txt(s, "WiFi Security Monitor", 1.5, 1.2, 10.3, 1.2, size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, "Your personal guard dog for your WiFi network", 1.5, 2.5, 10.3, 0.6, size=22, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "A simple guide — even a 10-year-old can follow!", 1.5, 3.2, 10.3, 0.5, size=16, color=MUTED, align=PP_ALIGN.CENTER)

# decorative cards at bottom
for i, (lbl, col) in enumerate([("Dashboard", ACCENT), ("Alerts", CRITICAL), ("Devices", HIGH)]):
    xl = 2.5 + i * 3
    box(s, xl, 5.0, 2.5, 1.2, fill_color=CARD)
    txt(s, lbl, xl, 5.0, 2.5, 1.2, size=18, bold=True, color=col, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2 — What is WiFi? (very simple)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

txt(s, "First — What is WiFi?", 0.5, 0.3, 12, 0.7, size=30, bold=True, color=ACCENT)
txt(s, "Think of WiFi like invisible radio waves that connect your phone, laptop,\nand TV to the internet — all through the air.", 0.5, 1.1, 12.3, 0.8, size=16, color=WHITE)

emoji_bullet(s, "📡", "Your router sends out WiFi signals like a radio station",             0.5, 2.1)
emoji_bullet(s, "📱", "Your devices (phone, laptop) listen and connect to those signals",    0.5, 2.65)
emoji_bullet(s, "🌍", "Once connected, you can browse the internet, watch videos, chat",     0.5, 3.2)
emoji_bullet(s, "🏠", "Everyone in your house shares the same WiFi network",                 0.5, 3.75)

box(s, 0.5, 4.6, 12.3, 1.5, fill_color=CARD)
txt(s, "⚠️  The Problem", 0.75, 4.65, 11, 0.35, size=13, bold=True, color=MEDIUM)
txt(s, "Because WiFi travels through the air, bad people nearby can try to spy on it,\nsteal your passwords, or even kick you off the internet — without touching your router!",
    0.75, 5.05, 11.8, 0.9, size=13, color=WHITE)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3 — The 4 attacks (simple language)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

txt(s, "4 Ways Bad Guys Attack Your WiFi", 0.5, 0.3, 12, 0.7, size=30, bold=True, color=ACCENT)
txt(s, "Our tool watches for all four of these, 24/7", 0.5, 1.0, 12, 0.4, size=15, color=MUTED)

attacks = [
    (CRITICAL, "👁  Evil Twin AP",      "MOST DANGEROUS",  "A hacker creates a fake WiFi network with the SAME name as yours.\nWhen your phone connects to it, they can see everything you type."),
    (HIGH,     "↔  ARP Spoofing",       "DANGEROUS",       "A hacker tricks your devices into sending internet traffic to THEM first.\nLike someone secretly reading your letters before delivering them."),
    (HIGH,     "⚡  Deauth Flood",      "DANGEROUS",       "A hacker keeps sending 'disconnect' messages to your devices.\nIt's like someone keep pressing mute on your TV remote over and over."),
    (MEDIUM,   "🔍  Port Scan",         "WARNING",         "A hacker is knocking on every 'door' of your network to find a weak spot.\nLike a burglar trying every window and door handle on your house."),
]

for i, (col, name, level, desc) in enumerate(attacks):
    row = i // 2
    col_x = 0.4 + (i % 2) * 6.4
    row_y = 1.7 + row * 2.5
    box(s, col_x, row_y, 6.1, 2.2, fill_color=CARD)
    box(s, col_x, row_y, 6.1, 0.05, fill_color=col)
    txt(s, name,  col_x+0.2, row_y+0.15, 4, 0.35, size=15, bold=True, color=col)
    txt(s, level, col_x+0.2, row_y+0.5,  5.7, 0.3, size=10, bold=True, color=col)
    txt(s, desc,  col_x+0.2, row_y+0.8,  5.7, 1.2, size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4 — What does this tool do?
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

txt(s, "What Does the WiFi Security Monitor Do?", 0.5, 0.3, 12.3, 0.7, size=28, bold=True, color=ACCENT)

box(s, 0.5, 1.2, 12.3, 1.0, fill_color=CARD)
txt(s, "🛡️  Think of it as a SECURITY GUARD for your WiFi network.", 0.75, 1.3, 11.8, 0.4, size=16, bold=True, color=WHITE)
txt(s, "It watches all the WiFi signals around you and immediately alerts you if something suspicious is happening.",
    0.75, 1.65, 11.8, 0.4, size=13, color=MUTED)

steps = [
    (ACCENT,    "1", "WATCHES",  "Listens to all WiFi signals floating through the air near your router, every single second"),
    (MEDIUM,    "2", "DETECTS",  "Automatically recognises the 4 types of attacks using smart detection rules and AI"),
    (CRITICAL,  "3", "ALERTS",   "Immediately shows a WARNING on the dashboard with full details about what happened"),
    (HIGH,      "4", "EXPLAINS", "Tells you exactly what the attack means and gives you step-by-step instructions to fix it"),
]

for i, (col, num, heading, body) in enumerate(steps):
    y = 2.45 + i * 1.0
    box(s, 0.5, y, 0.55, 0.55, fill_color=col)
    txt(s, num, 0.5, y, 0.55, 0.55, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, heading, 1.25, y+0.05, 2.5, 0.3, size=13, bold=True, color=col)
    txt(s, body,    3.9,  y+0.05, 9.0, 0.4, size=13, color=WHITE)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5 — The Dashboard
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

txt(s, "Page 1: The Dashboard", 0.5, 0.25, 8, 0.6, size=28, bold=True, color=ACCENT)
txt(s, "Your home screen — shows the big picture at a glance", 0.5, 0.85, 8, 0.4, size=14, color=MUTED)

# Left: explanation
items = [
    ("🔢", "Summary Numbers (top row)",   "Tells you at a glance: how many alerts are waiting, how many are critical, high, or medium. Like a scoreboard."),
    ("🎯", "Threat Meter (the gauge)",     "A big dial that goes from GREEN (safe) → ORANGE (watch out) → RED (danger!). Updated automatically every minute."),
    ("📊", "Alert Rate Chart (bar chart)", "Shows you WHEN attacks happened over the last hour. Tall bars = lots of attacks happening."),
    ("📋", "Live Feed (right side)",       "A live list of the 10 most recent alerts, newest at the top. Updates automatically — no need to refresh!"),
]

for i, (em, heading, body) in enumerate(items):
    y = 1.5 + i * 1.35
    txt(s, em,      0.5,  y,      0.5, 0.35, size=20)
    txt(s, heading, 1.1,  y,      7.2, 0.35, size=13, bold=True, color=WHITE)
    txt(s, body,    1.1,  y+0.38, 7.2, 0.7,  size=12, color=MUTED)

# Right: mock screen
box(s, 8.9, 0.25, 4.0, 7.0, fill_color=CARD)
txt(s, "DASHBOARD",           9.1, 0.4,  3.6, 0.35, size=10, bold=True, color=ACCENT)
# summary mini cards
for j, (val, col_) in enumerate([("1",CRITICAL),("2",HIGH),("1",MEDIUM),("3",ACCENT)]):
    bx = 9.1 + j * 0.9
    box(s, bx, 0.85, 0.82, 0.55, fill_color=BG)
    txt(s, val, bx, 0.85, 0.82, 0.55, size=18, bold=True, color=col_, align=PP_ALIGN.CENTER)
# gauge circle
txt(s, "⬤", 10.4, 1.65, 1.5, 1.5, size=60, color=HIGH, align=PP_ALIGN.CENTER)
txt(s, "65", 10.4, 1.9,  1.5, 0.6, size=22, bold=True, color=HIGH, align=PP_ALIGN.CENTER)
txt(s, "ELEVATED", 10.4, 2.5, 1.5, 0.3, size=9, color=MUTED, align=PP_ALIGN.CENTER)
# bars
for j, h in enumerate([0.2, 0.4, 0.15, 0.7, 0.3, 0.5]):
    box(s, 9.1+j*0.45, 4.3-h, 0.35, h, fill_color=HIGH)
txt(s, "Alert Rate", 9.1, 4.35, 3.6, 0.3, size=9, color=DIM)
# feed items
for j, (sev_lbl, col_) in enumerate([("CRIT",CRITICAL),("HIGH",HIGH),("HIGH",HIGH)]):
    box(s, 9.1, 4.8+j*0.55, 0.55, 0.3, fill_color=col_)
    txt(s, sev_lbl, 9.1, 4.8+j*0.55, 0.55, 0.3, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Alert detected...", 9.75, 4.83+j*0.55, 2.9, 0.3, size=9, color=MUTED)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6 — Alerts Page
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

txt(s, "Page 2: The Alerts Page", 0.5, 0.25, 8, 0.6, size=28, bold=True, color=ACCENT)
txt(s, "Where you investigate and take action on every alert", 0.5, 0.85, 8, 0.4, size=14, color=MUTED)

features = [
    ("🗂️", "Filter Bar (top of page)",     "Like a search filter! Choose to see only CRITICAL alerts, or only ARP Spoof attacks, or only unread ones. Tap any pill button to filter."),
    ("📋", "Alert Table (main list)",       "Every alert in a neat row: what severity it is, what type of attack, when it happened, which device was involved, and a short message."),
    ("🖱️", "Click any row → Detail Panel", "Clicking a row slides open a panel on the right with EVERYTHING: who the device is, what the attack means, and step-by-step instructions to fix it!"),
    ("✅", "Acknowledge Button",            "Once you have read and understood an alert, press 'Acknowledge' to mark it as done. Keeps your list clean."),
    ("🚨", "Burst Warning Banner",          "If 3 or more attacks of the same type happen within 2 minutes, a yellow banner appears at the top to warn you immediately."),
]

for i, (em, heading, body) in enumerate(features):
    y = 1.45 + i * 1.1
    txt(s, em,      0.5, y,      0.5, 0.35, size=18)
    txt(s, heading, 1.1, y,      7.2, 0.35, size=13, bold=True, color=WHITE)
    txt(s, body,    1.1, y+0.35, 7.2, 0.65, size=12, color=MUTED)

# Right: mock alerts table
box(s, 8.9, 0.25, 4.0, 7.0, fill_color=CARD)
txt(s, "ALERTS", 9.1, 0.4, 3.6, 0.3, size=10, bold=True, color=ACCENT)
# filter pills
for j, (lbl, act) in enumerate([("ALL",True),("CRIT",False),("HIGH",False)]):
    bx = 9.1 + j*1.2
    box(s, bx, 0.85, 1.1, 0.28, fill_color=(ACCENT if act else BG))
    txt(s, lbl, bx, 0.85, 1.1, 0.28, size=8, bold=act, color=(BG if act else MUTED), align=PP_ALIGN.CENTER)
# table header
for j, h in enumerate(["SEV","TYPE","TIME","ACTION"]):
    txt(s, h, 9.1+j*0.95, 1.3, 0.9, 0.25, size=8, bold=True, color=DIM)
# rows
rows = [
    (CRITICAL, "Evil Twin", "2m",  True),
    (HIGH,     "ARP Spoof", "15m", True),
    (HIGH,     "Deauth",    "31m", True),
    (MEDIUM,   "Port Scan", "47m", False),
]
for j, (col_, typ, tm, unacked) in enumerate(rows):
    y2 = 1.65 + j * 0.65
    if unacked:
        box(s, 9.05, y2-0.05, 3.9, 0.6, fill_color=RGBColor(0x1e,0x2d,0x45))
    box(s, 9.1, y2, 0.55, 0.3, fill_color=col_)
    txt(s, "●", 9.1, y2, 0.55, 0.3, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, typ, 9.75, y2+0.04, 0.9, 0.25, size=9, color=WHITE)
    txt(s, tm,  10.7, y2+0.04, 0.5, 0.25, size=9, color=MUTED)
    if unacked:
        box(s, 11.3, y2+0.03, 0.65, 0.25, fill_color=ACCENT)
        txt(s, "ACK", 11.3, y2+0.03, 0.65, 0.25, size=7, bold=True, color=BG, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7 — Devices Page
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

txt(s, "Page 3: The Devices Page", 0.5, 0.25, 8, 0.6, size=28, bold=True, color=ACCENT)
txt(s, "See every device on your network and its history", 0.5, 0.85, 8, 0.4, size=14, color=MUTED)

feats = [
    ("📦", "Device Cards",              "Every device on your network gets its own card showing: its name, IP address, MAC address, who made it (vendor), and when it was last seen."),
    ("🟢", "Active Indicator",          "Devices seen in the last 5 minutes glow with a CYAN border and show an 'Active' badge. Helps you spot who is currently online."),
    ("🏷️", "Alert Chips on cards",      "Each card shows coloured badges for any attacks linked to that device. E.g. '1× Critical' means that device was involved in 1 critical alert."),
    ("⚠️", "UNIDENTIFIED badge",        "If a device has no name and no known manufacturer — it shows an ORANGE 'UNIDENTIFIED' badge. This could be a rogue or suspicious device!"),
    ("📅", "Timeline (expandable)",     "Click 'Show timeline' on any card to see a horizontal line showing when the device first appeared and when any attacks happened on it."),
]

for i, (em, heading, body) in enumerate(feats):
    y = 1.45 + i * 1.1
    txt(s, em,      0.5, y,      0.5, 0.35, size=18)
    txt(s, heading, 1.1, y,      7.2, 0.35, size=13, bold=True, color=WHITE)
    txt(s, body,    1.1, y+0.35, 7.2, 0.65, size=12, color=MUTED)

# Right: mock device cards
box(s, 8.9, 0.25, 4.0, 3.3, fill_color=CARD, border_color=ACCENT)
txt(s, "johns-laptop",  9.1, 0.4,  3.0, 0.35, size=13, bold=True, color=WHITE)
txt(s, "aa:bb:cc:dd:ee:ff", 9.1, 0.75, 3.0, 0.25, size=9, color=DIM)
box(s, 12.1, 0.4, 0.65, 0.25, fill_color=ACCENT)
txt(s, "Active", 12.1, 0.4, 0.65, 0.25, size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)
txt(s, "IP   192.168.1.10", 9.1, 1.1, 3.6, 0.28, size=10, color=MUTED)
txt(s, "Vendor   Apple",    9.1, 1.38, 3.6, 0.28, size=10, color=MUTED)
box(s, 9.1, 1.75, 1.1, 0.28, fill_color=CRITICAL)
txt(s, "Critical ×1", 9.1, 1.75, 1.1, 0.28, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 10.3, 1.75, 0.9, 0.28, fill_color=HIGH)
txt(s, "High ×1", 10.3, 1.75, 0.9, 0.28, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# timeline
txt(s, "▼ Show timeline", 9.1, 2.15, 3.6, 0.25, size=9, color=DIM)
box(s, 9.1, 2.5, 3.6, 0.08, fill_color=DIM)
box(s, 9.5, 2.44, 0.12, 0.2, fill_color=CRITICAL)
box(s, 10.2, 2.44, 0.12, 0.2, fill_color=HIGH)

# Unidentified card
box(s, 8.9, 3.75, 4.0, 2.8, fill_color=CARD)
txt(s, "192.168.1.55",    9.1, 3.9,  3.0, 0.35, size=13, bold=True, color=MEDIUM)
txt(s, "de:ad:be:ef:00:55", 9.1, 4.25, 3.0, 0.25, size=9, color=DIM)
box(s, 9.1, 4.6, 1.3, 0.28, fill_color=RGBColor(0x3d,0x2e,0x00))
txt(s, "⚠ UNIDENTIFIED", 9.1, 4.6, 1.3, 0.28, size=8, bold=True, color=MEDIUM, align=PP_ALIGN.CENTER)
txt(s, "IP   192.168.1.55", 9.1, 5.0, 3.6, 0.28, size=10, color=MUTED)
box(s, 9.1, 5.4, 1.0, 0.28, fill_color=MEDIUM)
txt(s, "Medium ×1", 9.1, 5.4, 1.0, 0.28, size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8 — The Alert Detail Drawer
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

txt(s, "The Alert Detail Panel — Your Investigation Tool", 0.5, 0.25, 12, 0.6, size=26, bold=True, color=ACCENT)
txt(s, "Click any alert row → a panel slides in from the right with everything you need to understand and respond", 0.5, 0.9, 12, 0.4, size=13, color=MUTED)

sections = [
    (CRITICAL, "① Severity + Type badges",  "Shows you instantly how serious it is (Critical/High/Medium) and what kind of attack it is."),
    (HIGH,     "② Suspected Device info",   "Shows the device involved: its name, IP address, MAC address, manufacturer, and when it was first and last seen on your network."),
    (ACCENT,   "③ Alert Context",           "The raw alert message with full technical details, plus the unique alert ID so you can reference it later."),
    (MEDIUM,   "④ Attack Brief",            "A plain-English explanation of what the attack IS and WHY it is dangerous. Written so anyone can understand it."),
    (HIGH,     "⑤ Response Steps",          "A numbered list telling you EXACTLY what to do to stop the attack, in order of priority. No guesswork needed."),
    (MUTED,    "⑥ MITRE ATT&CK reference", "A link to the official global database of cyber attacks so you can learn even more if you want."),
]

for i, (col_, heading, body) in enumerate(sections):
    row = i // 2
    cx  = 0.5 + (i % 2) * 6.4
    cy  = 1.55 + row * 1.65
    box(s, cx, cy, 6.1, 1.5, fill_color=CARD)
    box(s, cx, cy, 6.1, 0.05, fill_color=col_)
    txt(s, heading, cx+0.18, cy+0.15, 5.7, 0.35, size=13, bold=True, color=col_)
    txt(s, body,    cx+0.18, cy+0.55, 5.7, 0.85, size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9 — How to use: step by step
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

txt(s, "How to Use the Tool — Step by Step", 0.5, 0.25, 12.3, 0.6, size=28, bold=True, color=ACCENT)
txt(s, "Follow these steps every time you want to check your network", 0.5, 0.9, 12.3, 0.4, size=14, color=MUTED)

steps2 = [
    (ACCENT,   "Open the Dashboard",    "Go to http://localhost:5173 in your browser. This is your home screen. Look at the Threat Meter first — is it green, orange, or red?"),
    (MEDIUM,   "Check the numbers",     "Look at the 4 boxes at the top. If any show a number other than zero — especially Critical — go to the Alerts page immediately."),
    (CRITICAL, "Go to Alerts page",     "Click 'Alerts' in the top menu. You'll see all recent alerts in a table. Unacknowledged alerts are highlighted."),
    (HIGH,     "Click an alert row",    "Click on any alert to open the detail panel. Read the Attack Brief and follow the Response Steps listed there."),
    (ACCENT,   "Acknowledge when done", "After you have dealt with an alert, click 'Acknowledge'. This marks it as handled and keeps your list clean."),
    (MUTED,    "Check Devices page",    "Go to the Devices page. Look for any ORANGE 'UNIDENTIFIED' badges — those devices need to be investigated."),
]

for i, (col_, heading, body) in enumerate(steps2):
    y = 1.55 + i * 0.97
    box(s, 0.5, y, 0.55, 0.55, fill_color=col_)
    txt(s, str(i+1), 0.5, y, 0.55, 0.55, size=18, bold=True, color=(BG if col_ != MUTED else WHITE), align=PP_ALIGN.CENTER)
    txt(s, heading, 1.25, y+0.06, 3.5, 0.3, size=13, bold=True, color=col_)
    txt(s, body,    5.0,  y+0.06, 7.8, 0.45, size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════
#  SLIDE 10 — Summary
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

box(s, 0, 3.1, 13.33, 0.04, fill_color=ACCENT)

txt(s, "Summary", 1.5, 0.5, 10.3, 0.8, size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, "What you have just learned", 1.5, 1.35, 10.3, 0.5, size=18, color=MUTED, align=PP_ALIGN.CENTER)

bullets = [
    "✅  WiFi signals travel through the air — anyone nearby could try to attack them",
    "✅  There are 4 main attack types: Evil Twin, ARP Spoof, Deauth Flood, Port Scan",
    "✅  The Dashboard shows your overall threat level with a live gauge and charts",
    "✅  The Alerts page lists every attack — click any row for full details and fix steps",
    "✅  The Devices page shows every device — watch out for UNIDENTIFIED ones",
    "✅  No technical knowledge needed — the tool explains everything in plain English",
]

for i, b in enumerate(bullets):
    txt(s, b, 1.5, 2.1 + i * 0.75, 10.3, 0.6, size=14, color=WHITE)

box(s, 1.5, 6.6, 10.3, 0.55, fill_color=CARD)
txt(s, "Open the prototype now:  double-click  dashboard.html  in the 'project 2 prototype' folder",
    1.7, 6.65, 9.9, 0.4, size=13, color=ACCENT)


# ── Save ───────────────────────────────────────────────────────
out = "c:/Users/oshiobughie/Desktop/project 2 prototype/WiFi_Security_Monitor_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
